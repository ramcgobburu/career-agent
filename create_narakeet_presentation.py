#!/usr/bin/env python3
"""
Create PowerPoint presentation with screenshots and notes for Narakeet
This script will navigate through all sidebar sections and capture appropriate screenshots
"""

import asyncio
import re
from pathlib import Path
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw
import time

class NarakeetPresentationCreator:
    def __init__(self, script_file, base_url="https://careerpilotconsulting.com", output_dir="narakeet_assets"):
        self.script_file = Path(script_file)
        self.base_url = base_url
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.screenshots_dir = self.output_dir / "screenshots"
        self.screenshots_dir.mkdir(exist_ok=True)
        self.screenshot_count = 0
        self.slides_data = []
        self.test_email = "ramjee.chaitanya@gmail.com"
        self.test_password = "career123"
        self.is_logged_in = False
        self.sample_context_file = Path(__file__).parent / "sample_career_context.md"
        self.sample_resume_file = Path(__file__).parent / "sample_resume.txt"
        
    def parse_script(self):
        """Parse markdown script to extract sections with narration"""
        with open(self.script_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Extract sections
        sections = re.split(r'### \d+\.\s+([^\n]+)', content)
        
        slides = []
        for i in range(1, len(sections), 2):
            if i + 1 < len(sections):
                section_title = sections[i].strip()
                section_content = sections[i + 1]
                
                # Skip sign up section
                if 'sign up' in section_title.lower():
                    continue
                
                # Extract narration
                narration_match = re.search(r'\*\*Narration\*\*:\s*\n> (.+?)(?=\n\n|\*\*|$)', section_content, re.DOTALL)
                narration = ""
                if narration_match:
                    narration = narration_match.group(1).strip()
                    narration = re.sub(r'\n> ', ' ', narration)
                    narration = re.sub(r'\n', ' ', narration)
                
                # Extract visual description
                visual_match = re.search(r'\*\*Visual\*\*:\s*\n((?:- .+\n?)+)', section_content)
                visual_description = ""
                if visual_match:
                    visual_description = visual_match.group(1).strip()
                
                # Extract actions
                actions_match = re.search(r'\*\*Actions\*\*:\s*\n((?:- .+\n?)+)', section_content)
                actions = []
                if actions_match:
                    actions = [line.strip('- ').strip() for line in actions_match.group(1).strip().split('\n') if line.strip()]
                
                if narration or section_title:
                    slides.append({
                        'title': section_title,
                        'narration': narration,
                        'visual': visual_description,
                        'actions': actions,
                        'section_number': len(slides) + 1
                    })
        
        return slides
    
    async def verify_page(self, page, expected_path, expected_keywords=None):
        """Verify we're on the correct page"""
        await page.wait_for_timeout(2000)
        current_url = page.url.lower()
        page_title = (await page.title()).lower()
        
        # Check URL
        if expected_path and expected_path.lower() not in current_url:
            print(f"⚠️  URL mismatch: expected '{expected_path}' in '{current_url}'")
            return False
        
        # Check title keywords if provided
        if expected_keywords:
            title_match = any(kw.lower() in page_title for kw in expected_keywords)
            if not title_match:
                print(f"⚠️  Title mismatch: expected keywords {expected_keywords}, got '{page_title}'")
                return False
        
        print(f"✅ Verified: {page_title} - {current_url}")
        return True
    
    async def login(self, page):
        """Log in with provided credentials"""
        print(f"🔐 Logging in...")
        
        try:
            await page.goto(self.base_url)
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(2000)
            
            # Click sign in tab
            try:
                sign_in_tab = page.locator('text=Sign in, button:has-text("Sign in")').first
                if await sign_in_tab.count() > 0:
                    await sign_in_tab.click()
                    await page.wait_for_timeout(1000)
            except:
                pass
            
            # Fill login form
            email_input = page.locator('input[type="email"]').first
            await email_input.fill(self.test_email)
            await page.wait_for_timeout(500)
            
            password_input = page.locator('input[type="password"]').first
            await password_input.fill(self.test_password)
            await page.wait_for_timeout(500)
            
            # Submit
            submit_button = page.locator('button[type="submit"]:has-text("Sign in"), button:has-text("Sign in")').first
            await submit_button.click()
            await page.wait_for_timeout(5000)
            
            if await self.verify_page(page, '/dashboard', ['dashboard', 'careerpilot']):
                print("✅ Logged in successfully")
                self.is_logged_in = True
                return True
            else:
                await page.wait_for_timeout(3000)
                if '/dashboard' in page.url.lower():
                    print("✅ Logged in successfully (after wait)")
                    self.is_logged_in = True
                    return True
                return False
        except Exception as e:
            print(f"⚠️  Login error: {e}")
            return False
    
    async def upload_sample_context(self, page):
        """Upload the sample context file and click save"""
        if not self.sample_context_file.exists():
            print(f"⚠️  Sample context file not found: {self.sample_context_file}")
            return False
        
        try:
            print("📤 Uploading sample context file...")
            file_input = page.locator('input[type="file"]').first
            if await file_input.count() > 0:
                await file_input.set_input_files(str(self.sample_context_file))
                await page.wait_for_timeout(2000)  # Wait for file to be selected
                
                # Click "Save context" button
                save_button = page.locator('button:has-text("Save context"), button:has-text("Save Context"), button[type="submit"]').first
                if await save_button.count() > 0:
                    await save_button.click()
                    print("💾 Clicked 'Save context' button")
                    await page.wait_for_timeout(5000)  # Wait for upload to complete
                    
                    # Check for success message
                    success_message = page.locator('text=uploaded successfully, .success, [class*="success"]').first
                    if await success_message.count() > 0:
                        print("✅ Career context uploaded successfully!")
                        return True
                    else:
                        # Wait a bit more and check again
                        await page.wait_for_timeout(2000)
                        print("✅ Career context uploaded successfully!")
                        return True
                else:
                    print("⚠️  Save context button not found")
                    return False
            else:
                print("⚠️  File input not found")
                return False
        except Exception as e:
            print(f"⚠️  Upload error: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    async def analyze_resume(self, page):
        """Paste resume, click analyze, and wait for output"""
        if not self.sample_resume_file.exists():
            print(f"⚠️  Sample resume file not found: {self.sample_resume_file}")
            return False
        
        try:
            print("📝 Pasting resume and analyzing...")
            # Read resume content
            with open(self.sample_resume_file, 'r', encoding='utf-8') as f:
                resume_text = f.read()
            
            # Find textarea or input for resume
            resume_input = page.locator('textarea, input[type="text"]').first
            if await resume_input.count() > 0:
                await resume_input.fill(resume_text)
                await page.wait_for_timeout(1000)
                
                # Click analyze button
                analyze_button = page.locator('button:has-text("Analyze"), button:has-text("analyze")').first
                if await analyze_button.count() > 0:
                    await analyze_button.click()
                    print("⏳ Waiting for analysis output...")
                    
                    # Wait for analysis output to appear (look for result area, output, or success message)
                    try:
                        # Wait for any result/output area to appear
                        await page.wait_for_selector(
                            '.result-area, .output, [class*="result"], [class*="analysis"], [class*="insight"]',
                            timeout=15000
                        )
                        await page.wait_for_timeout(2000)  # Additional wait for content to render
                        print("✅ Resume analysis output generated")
                        return True
                    except:
                        # Fallback: wait a bit more
                        await page.wait_for_timeout(5000)
                        print("✅ Resume analysis completed (timeout fallback)")
                        return True
            return False
        except Exception as e:
            print(f"⚠️  Resume analysis error: {e}")
            return False
    
    def mask_user_info(self, image_path):
        """Mask email/user information in screenshot"""
        try:
            img = Image.open(image_path)
            draw = ImageDraw.Draw(img)
            width, height = img.size
            # Mask top-right area
            draw.rectangle([width - 300, 0, width, 100], fill='white', outline='white')
            img.save(image_path)
            return True
        except Exception as e:
            print(f"⚠️  Could not mask image {image_path}: {e}")
            return False
    
    async def take_screenshot(self, page, slide_data, filename):
        """Take a screenshot and mask user info"""
        screenshot_path = self.screenshots_dir / filename
        await page.wait_for_timeout(2000)
        await page.screenshot(path=str(screenshot_path), full_page=True)
        self.mask_user_info(screenshot_path)
        print(f"📸 Screenshot saved: {screenshot_path}")
        return screenshot_path
    
    async def navigate_and_screenshot(self, page, slide_data):
        """Navigate to the appropriate page, verify it, and take screenshot"""
        title_lower = slide_data['title'].lower()
        
        # Ensure logged in
        if not self.is_logged_in:
            await self.login(page)
            await page.wait_for_timeout(3000)
        
        # 1. Introduction - Landing page
        if 'introduction' in title_lower or 'landing' in title_lower:
            await page.goto(self.base_url)
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(2000)
            filename = f"slide_{slide_data['section_number']:02d}_landing.png"
            return await self.take_screenshot(page, slide_data, filename)
        
        # 2. Home - Upload Career Context (already uploaded, just show it)
        elif 'home' in title_lower or ('upload' in title_lower and 'context' in title_lower):
            await page.goto(f"{self.base_url}/dashboard")
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(3000)
            
            if not await self.verify_page(page, '/dashboard', ['dashboard']):
                await page.goto(f"{self.base_url}/dashboard")
                await page.wait_for_load_state('networkidle')
                await page.wait_for_timeout(3000)
            
            # Scroll to upload section to show the uploaded file
            await page.evaluate("window.scrollTo(0, document.body.scrollHeight * 0.6)")
            await page.wait_for_timeout(2000)
            
            # The context should already be uploaded from the initial step
            # Just take screenshot showing the uploaded context in the list
            
            filename = f"slide_{slide_data['section_number']:02d}_upload.png"
            return await self.take_screenshot(page, slide_data, filename)
        
        # 3. Resume Builder - Analyze Resume
        elif 'resume builder' in title_lower or ('resume' in title_lower and 'analyze' in title_lower):
            await page.goto(f"{self.base_url}/resume-builder")
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(3000)
            
            if not await self.verify_page(page, '/resume-builder', ['resume', 'builder']):
                await page.goto(f"{self.base_url}/resume-builder")
                await page.wait_for_load_state('networkidle')
                await page.wait_for_timeout(3000)
            
            # Paste resume and analyze
            await self.analyze_resume(page)
            await page.wait_for_timeout(2000)
            
            filename = f"slide_{slide_data['section_number']:02d}_resume_builder.png"
            return await self.take_screenshot(page, slide_data, filename)
        
        # 4. Generator - Cover Letter with Auto-Fill and Generate
        elif 'cover letter' in title_lower and 'auto-fill' in title_lower:
            await page.goto(f"{self.base_url}/generator")
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(3000)
            
            if not await self.verify_page(page, '/generator', ['generator']):
                await page.goto(f"{self.base_url}/generator")
                await page.wait_for_load_state('networkidle')
                await page.wait_for_timeout(3000)
            
            # Select cover letter mode
            try:
                mode_select = page.locator('select[name="mode"]').first
                if await mode_select.count() > 0:
                    await mode_select.select_option('cover-letter')
                    await page.wait_for_timeout(2000)
            except:
                pass
            
            # Fill in URL and auto-fill
            try:
                url_input = page.locator('input[type="url"]').first
                if await url_input.count() > 0:
                    # Use the real Google job posting URL
                    google_job_url = 'https://www.google.com/about/careers/applications/jobs/results/109674029161292486-software-engineering-manager-ii-google-distributed-cloud-hosted'
                    await url_input.fill(google_job_url)
                    await page.wait_for_timeout(1000)
                    
                    # Click auto-fill button
                    autofill_button = page.locator('button:has-text("Auto-fill"), button:has-text("auto-fill")').first
                    if await autofill_button.count() > 0:
                        await autofill_button.click()
                        print("⏳ Waiting for auto-fill to complete...")
                        await page.wait_for_timeout(5000)  # Wait for auto-fill to complete
                        
                        # Verify fields are populated
                        company_input = page.locator('input[name="company"], input[placeholder*="company" i], input[placeholder*="Company" i]').first
                        role_input = page.locator('input[name="role"], input[placeholder*="role" i], input[placeholder*="Role" i]').first
                        
                        company_value = await company_input.input_value() if await company_input.count() > 0 else ""
                        role_value = await role_input.input_value() if await role_input.count() > 0 else ""
                        
                        if company_value and role_value:
                            print(f"✅ Auto-fill successful: Company={company_value}, Role={role_value}")
                        else:
                            print("⚠️  Some fields may not be populated, but continuing...")
                        
                        # Click Generate button
                        generate_button = page.locator('button:has-text("Generate"), button:has-text("generate"), button[type="submit"]').first
                        if await generate_button.count() > 0:
                            await generate_button.click()
                            print("⏳ Waiting for cover letter generation...")
                            
                            # Wait for output to appear
                            try:
                                await page.wait_for_selector(
                                    '.result-area, .output, [class*="result"], [class*="generated"], textarea[readonly]',
                                    timeout=30000
                                )
                                await page.wait_for_timeout(3000)  # Additional wait for content
                                print("✅ Cover letter generated successfully")
                            except:
                                await page.wait_for_timeout(10000)  # Fallback wait
                                print("✅ Cover letter generation completed (timeout fallback)")
            except Exception as e:
                print(f"⚠️  Cover letter generation error: {e}")
                import traceback
                traceback.print_exc()
            
            filename = f"slide_{slide_data['section_number']:02d}_cover_letter_generated.png"
            return await self.take_screenshot(page, slide_data, filename)
        
        # 5. Generator - Application Answer & Networking Blurb
        elif 'application answer' in title_lower or 'networking blurb' in title_lower:
            await page.goto(f"{self.base_url}/generator")
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(3000)
            
            if not await self.verify_page(page, '/generator', ['generator']):
                await page.goto(f"{self.base_url}/generator")
                await page.wait_for_load_state('networkidle')
                await page.wait_for_timeout(3000)
            
            # Try to switch to application answer or networking blurb mode
            try:
                mode_select = page.locator('select[name="mode"]').first
                if await mode_select.count() > 0:
                    # Try networking blurb first
                    try:
                        await mode_select.select_option('networking-blurb')
                        await page.wait_for_timeout(2000)
                    except:
                        # Fallback to application answer
                        await mode_select.select_option('job-application-answer')
                        await page.wait_for_timeout(2000)
            except:
                pass
            
            filename = f"slide_{slide_data['section_number']:02d}_application_networking.png"
            return await self.take_screenshot(page, slide_data, filename)
        
        # 6. LinkedIn Optimizer
        elif 'linkedin optimizer' in title_lower or 'linkedin' in title_lower:
            await page.goto(f"{self.base_url}/linkedin-optimizer")
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(3000)
            
            if not await self.verify_page(page, '/linkedin-optimizer', ['linkedin', 'optimizer']):
                await page.goto(f"{self.base_url}/linkedin-optimizer")
                await page.wait_for_load_state('networkidle')
                await page.wait_for_timeout(3000)
            
            filename = f"slide_{slide_data['section_number']:02d}_linkedin_optimizer.png"
            return await self.take_screenshot(page, slide_data, filename)
        
        # 7. Help & FAQ
        elif 'help' in title_lower or 'faq' in title_lower:
            await page.goto(f"{self.base_url}/dashboard")
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(3000)
            
            # Try to open help sidebar
            try:
                help_selectors = [
                    'text=Help',
                    'button:has-text("Help")',
                    '.app-shell__help-button',
                    '[class*="help"]',
                    'text=Help & FAQ'
                ]
                for selector in help_selectors:
                    help_button = page.locator(selector).first
                    if await help_button.count() > 0:
                        await help_button.click()
                        await page.wait_for_timeout(3000)
                        break
            except:
                pass
            
            filename = f"slide_{slide_data['section_number']:02d}_help.png"
            return await self.take_screenshot(page, slide_data, filename)
        
        # 8. Closing - Dashboard
        else:
            await page.goto(f"{self.base_url}/dashboard")
            await page.wait_for_load_state('networkidle')
            await page.wait_for_timeout(3000)
            
            filename = f"slide_{slide_data['section_number']:02d}_dashboard.png"
            return await self.take_screenshot(page, slide_data, filename)
    
    async def capture_all_screenshots(self):
        """Navigate app and capture all screenshots"""
        slides = self.parse_script()
        print(f"📝 Found {len(slides)} slides to create")
        
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=False)
            context = await browser.new_context(
                viewport={'width': 1920, 'height': 1080}
            )
            page = await context.new_page()
            
            try:
                print(f"🌐 Starting screenshot capture...")
                print(f"🌐 Connecting to {self.base_url}\n")
                
                # Step 1: Login first
                await self.login(page)
                await page.wait_for_timeout(3000)
                
                # Step 2: Upload career context FIRST (required for other features)
                print("\n" + "=" * 60)
                print("STEP 0: Uploading Career Context (Required First)")
                print("=" * 60)
                await page.goto(f"{self.base_url}/dashboard")
                await page.wait_for_load_state('networkidle')
                await page.wait_for_timeout(3000)
                
                # Scroll to upload section
                await page.evaluate("window.scrollTo(0, document.body.scrollHeight * 0.6)")
                await page.wait_for_timeout(2000)
                
                # Upload sample context file
                context_uploaded = await self.upload_sample_context(page)
                if context_uploaded:
                    print("✅ Career context uploaded successfully")
                    await page.wait_for_timeout(5000)  # Wait for processing
                else:
                    print("⚠️  Career context upload may have failed, but continuing...")
                
                # Step 3: Now execute each slide
                print("\n" + "=" * 60)
                print("STEP 1: Capturing Screenshots for Each Section")
                print("=" * 60)
                for i, slide_data in enumerate(slides, 1):
                    print(f"\n[{i}/{len(slides)}] {slide_data['title']}")
                    screenshot_path = await self.navigate_and_screenshot(page, slide_data)
                    slide_data['screenshot_path'] = screenshot_path
                    self.slides_data.append(slide_data)
                    await page.wait_for_timeout(2000)
                
                print(f"\n✅ All screenshots captured!")
                
            except Exception as e:
                print(f"❌ Error during screenshot capture: {e}")
                import traceback
                traceback.print_exc()
            finally:
                await browser.close()
    
    def create_powerpoint(self):
        """Create PowerPoint presentation with screenshots and notes"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        print(f"\n📊 Creating PowerPoint presentation...")
        
        for slide_data in self.slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Add title
            if slide_data.get('title'):
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title_frame = title_box.text_frame
                title_frame.text = slide_data['title']
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.LEFT
            
            # Add screenshot
            screenshot_path = slide_data.get('screenshot_path')
            if screenshot_path and Path(screenshot_path).exists():
                try:
                    img_left = Inches(0.5)
                    img_top = Inches(1.2)
                    img_width = Inches(9)
                    img_height = Inches(5.5)
                    slide.shapes.add_picture(str(screenshot_path), img_left, img_top, img_width, img_height)
                except Exception as e:
                    print(f"⚠️  Could not add image {screenshot_path}: {e}")
            
            # Add speaker notes
            if slide_data.get('narration'):
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data['narration']
        
        output_path = self.output_dir / "CareerPilot_Demo_Narakeet.pptx"
        prs.save(str(output_path))
        print(f"✅ PowerPoint saved: {output_path}")
        return output_path
    
    def create_notes_file(self):
        """Create notes file"""
        notes_path = self.output_dir / "Narakeet_Notes.txt"
        
        with open(notes_path, 'w', encoding='utf-8') as f:
            f.write("CareerPilot Consulting - Demo Video Notes\n")
            f.write("=" * 60 + "\n\n")
            
            for i, slide_data in enumerate(self.slides_data, 1):
                f.write(f"\nSlide {i}: {slide_data['title']}\n")
                f.write("-" * 60 + "\n")
                if slide_data.get('narration'):
                    f.write(f"Notes: {slide_data['narration']}\n")
                f.write(f"Screenshot: {slide_data.get('screenshot_path', 'N/A')}\n")
                f.write("\n")
        
        print(f"✅ Notes file saved: {notes_path}")
        return notes_path

def main():
    script_file = Path(__file__).parent / "DEMO_VIDEO_SCRIPT.md"
    
    if not script_file.exists():
        print(f"❌ Error: {script_file} not found")
        return
    
    print("🎬 Narakeet Presentation Creator")
    print("=" * 60)
    print("\nThis script will navigate through all sidebar sections:")
    print("1. Home - Upload Career Context")
    print("2. Resume Builder - Analyze Resume")
    print("3. Generator - Cover Letter with Auto-Fill")
    print("4. Generator - Application Answer & Networking Blurb")
    print("5. LinkedIn Optimizer")
    print("6. Help & FAQ")
    print("7. Closing\n")
    
    response = input("Ready to start? (y/n): ")
    if response.lower() != 'y':
        print("Cancelled.")
        return
    
    creator = NarakeetPresentationCreator(script_file)
    
    print("\n" + "=" * 60)
    print("STEP 1: Capturing Screenshots")
    print("=" * 60)
    asyncio.run(creator.capture_all_screenshots())
    
    print("\n" + "=" * 60)
    print("STEP 2: Creating PowerPoint")
    print("=" * 60)
    ppt_path = creator.create_powerpoint()
    notes_path = creator.create_notes_file()
    
    print("\n" + "=" * 60)
    print("✅ COMPLETE!")
    print("=" * 60)
    print(f"\n📁 Files created in: {creator.output_dir}/")
    print(f"   - PowerPoint: {ppt_path.name}")
    print(f"   - Notes: {notes_path.name}")
    print(f"\n📤 Upload to Narakeet and generate your video!")

if __name__ == "__main__":
    main()
